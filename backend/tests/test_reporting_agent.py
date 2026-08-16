import os
import sys
import zipfile
import io
import uuid
import datetime
import pytest
from sqlalchemy import create_engine, text
from sqlalchemy.orm import sessionmaker
from dotenv import load_dotenv

import openpyxl
from pptx import Presentation

sys.path.append(os.path.dirname(os.path.dirname(__file__)))
from agents.reporting import generate_executive_report
from models import AuditLog, ModelTrustScore
from agents.audit import compute_hash

load_dotenv(os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), '.env'))
db_url = os.getenv("DATABASE_URL")
engine = create_engine(db_url)
SessionLocal = sessionmaker(bind=engine)

def seed_reporting_data(session_id):
    db = SessionLocal()
    # Clean previous test runs for this session_id just in case
    db.query(AuditLog).filter(AuditLog.session_id == session_id).delete()
    
    # 1. Seed ModelTrustScore
    score = ModelTrustScore(
        id=str(uuid.uuid4()),
        timestamp=datetime.datetime.utcnow(),
        precision=0.95,
        recall=0.88,
        false_positive_rate=0.02,
        sample_size=1000
    )
    db.add(score)
    
    # 2. Seed some flagged transactions in AuditLog
    # Need to simulate the output of the investigation graph
    payload_flagged = {"amount": 25000, "merchant_category": "Crypto"}
    result_flagged = {"recommended_action": "FLAG", "risk_score": 92.5}
    
    # Simple hash for test
    h1 = compute_hash("test-exec-1", session_id, "investigation", "investigate_node", "review", payload_flagged, result_flagged, 100, 0.01, "Investigate this.", "Flagged for manual review.", "GENESIS")
    
    log1 = AuditLog(
        id=str(uuid.uuid4()),
        timestamp=datetime.datetime.utcnow(),
        execution_id="test-exec-1",
        session_id=session_id,
        record_type="investigation",
        node_name="investigate_node",
        action_type="review",
        payload=payload_flagged,
        result=result_flagged,
        response="Flagged for manual review.",
        prev_row_hash="GENESIS",
        current_hash=h1
    )
    db.add(log1)
    
    # 3. Seed a variance query in AuditLog
    payload_var = {"sql": "SELECT SUM(amount)..."}
    result_var = {"final_answer": "Software expenses spiked due to VendorC."}
    
    h2 = compute_hash("test-exec-2", session_id, "variance_query", "sql_agent", "variance_analysis", payload_var, result_var, 50, 0.005, "Why did expenses spike?", "Software expenses spiked due to VendorC.", h1)
    
    log2 = AuditLog(
        id=str(uuid.uuid4()),
        timestamp=datetime.datetime.utcnow(),
        execution_id="test-exec-2",
        session_id=session_id,
        record_type="variance_query",
        node_name="sql_agent",
        action_type="variance_analysis",
        payload=payload_var,
        result=result_var,
        prompt="Why did expenses spike?",
        response="Software expenses spiked due to VendorC.",
        prev_row_hash=h1,
        current_hash=h2
    )
    db.add(log2)
    
    db.commit()
    db.close()


def test_executive_report_structure():
    session_id = f"test-session-{uuid.uuid4()}"
    seed_reporting_data(session_id)
    
    # 1. Call the agent to generate the report bundle
    zip_bytes = generate_executive_report(session_id)
    
    assert isinstance(zip_bytes, bytes), "Agent must return raw bytes for the zip archive"
    
    # 2. Open the zip and verify structure
    with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
        file_list = zf.namelist()
        assert "executive_report.pptx" in file_list
        assert "investigation_details.xlsx" in file_list
        
        pptx_data = zf.read("executive_report.pptx")
        xlsx_data = zf.read("investigation_details.xlsx")
        
    # 3. Verify XLSX Structure
    wb = openpyxl.load_workbook(io.BytesIO(xlsx_data))
    sheet_names = wb.sheetnames
    assert "Case Details" in sheet_names
    assert "Trust Score Trend" in sheet_names
    
    # 4. Verify PPTX Structure
    prs = Presentation(io.BytesIO(pptx_data))
    assert len(prs.slides) == 3, "Presentation must have exactly 3 slides (Summary, Highlights, Actions)"
    
    # Slide 1 check
    slide_1_text = ""
    for shape in prs.slides[0].shapes:
        if hasattr(shape, "text"):
            slide_1_text += shape.text
            
    assert "Executive Summary" in slide_1_text or "Trust Score" in slide_1_text
